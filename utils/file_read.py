import PyPDF2
from flask import jsonify
from pptx import Presentation
from docx import Document


def process_file(path, filetype):
    try:
        if filetype == 'pdf':
            return extract_text_from_pdf(path)
        elif filetype == 'pptx':
            return extract_text_from_ppt(path)
        elif filetype == 'txt':
            return extract_text_from_txt(path) 
        elif filetype == 'docx':
            return extract_text_from_doc(path)
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
def extract_text_from_pdf(filepath):
    # Opening the PDF file in binary read mode
    with open(filepath, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + " "
    return text

def extract_text_from_txt(filepath):
    with open(filepath, 'r', encoding='utf-8') as file:
        text = file.read()
    return text

def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_doc(path):
    doc = Document(path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text
