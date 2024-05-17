import io
import json
import os
import re
import time
from datetime import datetime 
import math
from concurrent.futures import ThreadPoolExecutor
import PyPDF2
import requests
import tiktoken
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from flask import Flask, jsonify, request, g
from flask_cors import CORS
from gevent import monkey
from openai import OpenAI
from werkzeug.utils import secure_filename
from pymongo import MongoClient
from youtube_transcript_api import YouTubeTranscriptApi

load_dotenv()

api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=api_key)

app = Flask(__name__)
CORS(app)

MODEL = "gpt-3.5-turbo"
UPLOAD_FOLDER = "uploads"

SUMMARY_PROMPT_FILE_PATH = 'Prompts/summary.txt'
QUIZ_PROMPT_FILE_PATH = 'Prompts/quiz.txt'
STUB_PROMPT_FILE_PATH = 'Prompts/stub.txt'

MAX_RETRIES = 5  
RETRY_DELAY = 1

mongo_client = MongoClient('mongodb+srv://krish:yXMdTPwSdTRo7qHY@serverlessinstance0.18otqeg.mongodb.net/')
db = mongo_client['Lurny']
collection = db['contents']

def get_transcript(video_id):
    """Function get_transcript from youtube video id"""        
    transcript = YouTubeTranscriptApi.get_transcript(video_id)

def is_youtube_url(url):
    """Function check url is youtube url"""    
    
    return "youtube.com" in url or "youtu.be" in url

def num_tokens_from_string(string: str, encoding_name: str) -> int:
    """Function get number of tokens"""
    encoding = tiktoken.encoding_for_model(encoding_name)
    num_tokens = len(encoding.encode(string))
    return num_tokens

def split_content_evenly(content, parts):
    """Function split content as parts"""
    
    if parts <= 0:
        raise ValueError("Number of parts must be greater than zero.")
    
    part_len = len(content) // parts
    splits = []
    index = 0
    
    for _ in range(parts):
        splits.append(content[index: index + part_len])
        index += part_len
    return splits

def summarize(text):
    """Function summarize with retry logic"""    
    
    tools = [
        {
            "type": "function",
            "function": {
                "name": "get_title_summary_hashtags",
                "description": "Get title, summary and hash_tags from use's message.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "title",
                        },
                        "summary": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "bulleted summary",
                        },
                        "hash_tags": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "hash tag",
                        },
                    },
                    "required": ["title", "summary", "hash_tags"],
                },
            },
        }
    ]

    attempts = 0

    while attempts < MAX_RETRIES:
        try:
            with open(SUMMARY_PROMPT_FILE_PATH, encoding='utf-8') as file:
                prompt = file.read()
            print("input token", num_tokens_from_string(prompt + str(tools), MODEL))
            response = client.chat.completions.create(
                model=MODEL,
                messages=[
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": text},
                ],
                tools=tools,
            )
            
            # Check if the response contains the expected output
            if response.choices[0].message.tool_calls is not None:
                output = []
                for res in response.choices[0].message.tool_calls:
                    output.append(res.function.arguments)
                return output[0]
            
            # If we got a response but tool_calls is None, raise an exception to retry
            raise ValueError("Received None from tool_calls")

        except (ValueError, AttributeError) as e:
            attempts += 1
            print(f"Attempt {attempts} failed (Summarize) with error: {e}. Retrying...")
            time.sleep(RETRY_DELAY)
    
    # If we've exceeded the maximum number of retries, raise an exception
    raise RuntimeError("Max retries exceeded. Unable to get a valid response.")

def quiz_from_stub(text):
    """Function Generate Quizzes"""
    # Assuming split_content_evenly is a function that splits the text into even parts
    splits = split_content_evenly(text, 5)
    
    quizzes = []
    # Use ThreadPoolExecutor to execute tasks asynchronously
    with ThreadPoolExecutor() as executor:
        # Map the generate_quiz function over the splits
        future_quizzes = executor.map(quiz, splits)
        
        for future_quiz in future_quizzes:
            quizzes.append(future_quiz)
    
    return quizzes   

def quiz(text):
    """Function Generate Quiz"""

    tools = [
        {
            "type": "function",
            "function": {
                "name": "get_question_answer",
                "description": "Generate question and answers, correctanswer, explanation for the question",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "question": {
                            "type": "string",
                            "description": "question",
                        },
                        "answer": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "answers",
                        },
                        "correctanswer": {
                            "type": "string",
                            "description": "correct answer",
                        },
                        "explanation": {"type": "string", "description": "explanation"},
                    },
                    "required": ["question", "answer", "correctanswer", "explanation"],
                },
            },
        }
    ]
    print("token number of split", num_tokens_from_string(text, MODEL))
    attempts = 0

    while attempts < MAX_RETRIES:
        try:
            with open(QUIZ_PROMPT_FILE_PATH, encoding='utf-8') as file:
                prompt = file.read()

            response = client.chat.completions.create(
                model=MODEL,
                messages=[
                    {"role": "system", "content": prompt},
                    {"role": "user", "content": text},
                ],
                tools=tools,
            )

            if response.choices[0].message.tool_calls is not None:
                output = []
                for res in response.choices[0].message.tool_calls:
                    output.append(res.function.arguments)
                return output[0]

            # If we got a response but tool_calls is None, raise an exception to retry
            raise ValueError("Received None from tool_calls")

        except (ValueError, AttributeError) as e:
            attempts += 1
            print(f"Attempt {attempts} failed (Quiz) with error: {e}. Retrying...")
            time.sleep(RETRY_DELAY)

    # If we've exceeded the maximum number of retries, raise an exception
    raise RuntimeError("Max retries exceeded. Unable to get a valid response.")

def extract_cover_image(url):
    """Function extract_cover_image"""    
    
    page = requests.get(url)
    soup = BeautifulSoup(page.content, "html.parser")
    cover_image = soup.find(
        "meta", property="og:image"
    )  # Check for Open Graph image tag
    if cover_image:
        return cover_image["content"]
    cover_image = soup.find("meta", itemprop="image")  # Check for Schema.org image tag
    if cover_image:
        return cover_image["content"]
    cover_image = soup.find("img")  # Fallback to the first image tag on the page
    if cover_image:
        return cover_image["src"]
    return None

def save_content(url, content):
  """Function save content to database"""
  document = {
    'url': url,
    'content': content,
    'timestamp': datetime.utcnow()
  }
  collection.insert_one(document)

@app.route("/manually", methods=["POST"])
def lurnify_from_content():
    """Function fetch_data_from_url"""       
    try:
        data = request.get_json()
        combined_text = data["content"]

        media = None
        url = None
        cover_image_url = None

        start = time.time()

        # Save content to DB
        save_content(None, combined_text)
        
        # Calulate number of tokens for the certain gpt model
        num_tokens = num_tokens_from_string(combined_text, MODEL)
        print(num_tokens, len(combined_text))
        
        num_parts = math.ceil(num_tokens / 15000) 
        # Split entire content to several same parts when content is large than gpt token limit
        splits = split_content_evenly(combined_text, num_parts) 
                
        results = []
        for split in splits:                        
            summary_content = summarize(split)
            question_content = quiz_from_stub(split)

            json_string = json.dumps(
                {
                    "summary_content": summary_content,
                    "questions": question_content,
                    "image": cover_image_url if cover_image_url is not None else "",
                    "url": url,
                    "media": media,
                }
            )
            results.append(json_string)
        print(results)

        end = time.time()
        print(end - start, "s")

        return jsonify(results)

    except Exception as e:
        return jsonify({"error": str(e)}), 500
 
@app.route("/fetch", methods=["POST"])
def lurnify_from_url():
    """Function lurnify_from_url"""       
    try:
        data = request.get_json()
        url = data["url"]

        # Initialize variables for the extracted text and media type
        combined_text = ""
        media = None
        cover_image_url = None

        start = time.time()

        # openai_client = g.openai_client
        # Extract text from the URL
        if is_youtube_url(url):
            # Extract video ID from the YouTube URL
            video_id_match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11})", url)
            if video_id_match:
                video_id = video_id_match.group(1)
                try:
                    # Get the transcript from the YouTube video
                    transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
                    combined_text = "".join(entry["text"] for entry in transcript_list)
                    print("Video script", combined_text)
                except Exception as e:
                    return jsonify({"error": str(e)}), 400
            else:
                return jsonify({"error": "No YouTube video ID found in URL"}), 400

        elif url.lower().endswith(".pdf"):
            media = "PDF"
            try:
                response = requests.get(url)
                response.raise_for_status()  # Ensure we capture HTTP errors

                pdf_reader = PyPDF2.PdfReader(io.BytesIO(response.content))
                combined_text = "".join(
                    page.extract_text() + " " for page in pdf_reader.pages
                )

            except Exception as e:
                return jsonify({"error": str(e)}), 400
        else:
            media = "web"
            page = requests.get(url)

            soup = BeautifulSoup(page.content, "html.parser")

            cover_image_url = extract_cover_image(url)

            text_elements = [
                tag.get_text() for tag in soup.find_all(["p", "span", "a", "li"])
            ]
            combined_text = "".join(text_elements).strip()
        
        # Save content to DB
        save_content(url, combined_text)
        
        # Calulate number of tokens for the certain gpt model
        num_tokens = num_tokens_from_string(combined_text, MODEL)
        print(num_tokens, len(combined_text))
        
        num_parts = math.ceil(num_tokens / 15000) 
        # Split entire content to several same parts when content is large than gpt token limit
        splits = split_content_evenly(combined_text, num_parts) 
                
        results = []
        for split in splits:                        
            summary_content = summarize(split)
            question_content = quiz_from_stub(split)

            json_string = json.dumps(
                {
                    "summary_content": summary_content,
                    "questions": question_content,
                    "image": cover_image_url if cover_image_url is not None else "",
                    "url": url,
                    "media": media,
                }
            )
            results.append(json_string)
        print(results)

        end = time.time()
        print(end - start, "s")

        return jsonify(results)

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/lurnify-from-file', methods=['POST'])
def lurnify_from_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join("uploads", filename)
        file.save(file_path)
        text = process_file(file_path, filename.rsplit('.', 1)[1].lower())
        print("text", text)
        start = time.time()

        # Save content to DB
        # save_content(None, text)
        
        # Calulate number of tokens for the certain gpt model
        num_tokens = num_tokens_from_string(text, MODEL)
        print(num_tokens, len(text))
        
        num_parts = math.ceil(num_tokens / 15000) 
        # Split entire content to several same parts when content is large than gpt token limit
        splits = split_content_evenly(text, num_parts) 
                
        results = []
        for split in splits:                        
            summary_content = summarize(split)
            question_content = quiz_from_stub(split)

            json_string = json.dumps(
                {
                    "summary_content": summary_content,
                    "questions": question_content,
                    "image": None,
                    "url": filename,
                    "media": "file",
                }
            )
            results.append(json_string)
        print(results)

        end = time.time()
        print(end - start, "s")

        return jsonify(results)
 
    else:
        return jsonify({'error': 'File type not allowed'}), 400

def process_file(path, filetype):
    try:
        if filetype == 'pdf':
            return extract_text_from_pdf(path)
        elif filetype == 'pptx':
            return extract_text_from_ppt(path)
        elif filetype == 'docx':
            return extract_text_from_doc(path)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def extract_text_from_pdf(path):
    with open(path, 'rb') as file:
        reader = PyPDF2.PdfFileReader(file)
        text = ''
        for page_num in range(reader.numPages):
            text += reader.getPage(page_num).extractText()
        return text

def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_doc(path):
    doc = Document(path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text

@app.route("/get_quiz", methods=["POST"])
def generage_quiz():
    """Function analyze"""    
    data = request.get_json()
    stub = data["stub"]

    tools = [
        {
            "type": "function",
            "function": {
                "name": "get_question_answer",
                "description": "Generate question and answers",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "question": {
                            "type": "string",
                            "description": "question",
                        },
                        "answer": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "available answer",
                        },
                        "correctanswer": {
                            "type": "string",
                            "description": "correct answer",
                        },
                        "explanation": {"type": "string", "description": "explanation"},
                    },
                    "required": ["question", "answer", "correctanswer", "explanation"],
                },
            },
        }
    ]    
    
    with open(STUB_PROMPT_FILE_PATH, encoding='utf-8') as file:
        prompt = file.read()
        
    response = client.chat.completions.create(
        model=MODEL,
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": stub},
        ],
        tools=tools,
    )    

    output = []
    for res in response.choices[0].message.tool_calls:
        output.append(res.function.arguments)
        print("output", output)
    return output[0] 

@app.route("/update_prompts", methods=["POST"])
def update_prompts():
    """Endpoint to update prompt files based on JSON input"""
    data = request.get_json()
    
    summary_prompt = data.get('summary')
    quiz_prompt = data.get('quiz')
    stub_prompt = data.get('stub')
    
    # Check if any of the prompts are provided and update accordingly
    if summary_prompt and not update_prompt_file(SUMMARY_PROMPT_FILE_PATH, summary_prompt):
        return jsonify({"error": f"Failed to update summary prompt"}), 500
    
    if quiz_prompt and not update_prompt_file(QUIZ_PROMPT_FILE_PATH, quiz_prompt):
        return jsonify({"error": "Failed to update quiz prompt"}), 500
    
    if stub_prompt and not update_prompt_file(STUB_PROMPT_FILE_PATH, stub_prompt):
        return jsonify({"error": "Failed to update stub prompt"}), 500
    
    return jsonify({"message": "Prompts updated successfully"}), 200

@app.route("/get_prompts", methods=["GET"])
def get_prompts():
    """Endpoint to get the current prompts"""
    summary_prompt = read_prompt_file(SUMMARY_PROMPT_FILE_PATH)
    quiz_prompt = read_prompt_file(QUIZ_PROMPT_FILE_PATH)
    stub_prompt = read_prompt_file(STUB_PROMPT_FILE_PATH)

    # Ensure all prompts were read successfully
    if summary_prompt is None or quiz_prompt is None or stub_prompt is None:
        return jsonify({"error": "Failed to read one or more prompt files"}), 500
    
    # Return the prompts in JSON format
    prompts = {
        "summary": summary_prompt,
        "quiz": quiz_prompt,
        "stub": stub_prompt,
    }
    
    return jsonify(prompts), 200

def read_prompt_file(file_path):
    """Function to read the content of a prompt file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return None
          
def update_prompt_file(file_path, new_prompt):
    """Function to update a prompt file"""
    try:
        with open(file_path, 'w', encoding='utf-8') as file:
            file.write(new_prompt)
        return True
    except Exception as e:
        return False

def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx'}
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

if __name__ == "__main__":
    app.run(
        host="0.0.0.0",
        port=5173,
        threaded=True,
        debug=True,
        use_reloader=False,
        ssl_context=(
            "/etc/letsencrypt/live/lurny.net/cert.pem",
            "/etc/letsencrypt/live/lurny.net/privkey.pem",
        ),
    )
